"""
One-shot PowerPoint generator for the "Modding Timberborn" tech-guild
talk that this whole project was built for.

Produces ./timberborn-talk.pptx — a ~36-slide deck targeting ~55
minutes of content + 10 minutes of Q&A. Each slide carries speaker
notes that show up in PowerPoint's presenter view.

Re-run with `python talk/generate.py` whenever the content changes.
Output is checked in so the projector laptop doesn't need Python or
python-pptx.

Talk arc:
  1. Why this matters here — Kallikor's sim platform vs Timberborn's.
  2. The customer (Jasper, age 6) and his three feature requests.
  3. Timberborn the game + how its mod loader works.
  4. Three layers of modding: templates → exposed interfaces → Harmony.
  5. Mod #1 (Quadruple Platform) — pure JSON + a Python-generated
     .timbermesh model. Templates layer.
  6. Mod #2 (Flood Season) — flow multiplier (interfaces) + flood as
     a new hazard (Harmony). Plus the cascade of vanilla code that
     fights back, the save/load story, and the fake-event knock-ons.
  7. Mod #3 (Mixed Tide) — doing it twice. Same patches, parallel
     branches. New gotchas surface.
  8. Mod #4 (Rainy Season, future) — what we'd need to find to make
     water fall on the map.
  9. Wrap, thanks, Q&A.

Library: python-pptx (pip install python-pptx).
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


# --- Personalisation ---------------------------------------------------------
# Fill these in before running the script for delivery.

SON_NAME = "Jasper"
SON_AGE = 6
PRESENTER = "Chris Brett"
TALK_DATE = "2026-09-10"
TALK_VENUE = "Kallikor tech guild"
GITHUB_URL = "https://github.com/Spycho/timberborn-flood-mod"

# Jasper's actual feature requests — replace if any of these are
# off and you want exact wording on the slides. The "too much water"
# request and the "flood every now and then" request were two distinct
# asks that landed in the same Flood Season mod as two phases (more
# flow during temperate, then flood as a new hazard).
JASPER_PLATFORM_BRIEF = "I want a platform that goes higher than three."
JASPER_TEMPERATE_QUOTE = "Dad, what if there was too much water?"
JASPER_HAZARD_QUOTE = "Dad, what if there was a flood every now and then, like the badtides and droughts?"
JASPER_MIXED_QUOTE = "Dad, let's make a season with bad and good water."
JASPER_RAIN_BRIEF = "What if water fell from the sky and made the map wet?"

# Closing slide line. Replace with something specific you actually
# learnt. Placeholder makes the deck read end-to-end during iteration.
CLOSING_REFLECTION = (
    "TODO: one or two sentences on what pair-programming with a "
    "kindergartener actually taught you. Best PM I've ever had / "
    "customer-driven scoping / something true and specific. "
    "Script ships with a placeholder."
)


# --- Theme -------------------------------------------------------------------

NAVY = RGBColor(0x0A, 0x20, 0x3A)
DEEP_BLUE = RGBColor(0x12, 0x40, 0x6E)
MID_BLUE = RGBColor(0x2C, 0x74, 0xAF)
ACCENT = RGBColor(0xAA, 0xDC, 0xF0)
OFFWHITE = RGBColor(0xF4, 0xF7, 0xFB)
DARK_TEXT = RGBColor(0x10, 0x18, 0x24)
MUTED = RGBColor(0x55, 0x60, 0x70)
WARN_RED = RGBColor(0xB8, 0x33, 0x33)
CODE_BG = RGBColor(0xF2, 0xF4, 0xF7)
CODE_TEXT = RGBColor(0x10, 0x18, 0x24)
COMPARE_LEFT_BG = RGBColor(0xE9, 0xF1, 0xF9)
COMPARE_RIGHT_BG = RGBColor(0xF8, 0xEE, 0xE9)


# --- Geometry ---------------------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

TITLE_TOP = Inches(0.35)
TITLE_HEIGHT = Inches(0.9)
BODY_TOP = Inches(1.4)
BODY_HEIGHT = Inches(5.6)
MARGIN_X = Inches(0.6)
BODY_WIDTH = SLIDE_W - MARGIN_X * 2
FOOTER_TOP = Inches(7.1)
FOOTER_HEIGHT = Inches(0.35)


# --- Low-level slide builders -----------------------------------------------


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_background(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height,
                 text="", *, font="Calibri", size=20, bold=False,
                 color=DARK_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if text:
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tf


def _add_title(slide, title: str):
    _add_textbox(
        slide,
        MARGIN_X, TITLE_TOP, SLIDE_W - MARGIN_X * 2, TITLE_HEIGHT,
        text=title, font="Calibri", size=32, bold=True,
        color=NAVY, anchor=MSO_ANCHOR.MIDDLE,
    )
    line = slide.shapes.add_connector(
        1, MARGIN_X, Inches(1.25), SLIDE_W - MARGIN_X, Inches(1.25),
    )
    line.line.color.rgb = MID_BLUE
    line.line.width = Pt(1.5)


def _add_footer(slide, slide_num: int, total: int, section: str = ""):
    text = f"{section}   ·   {slide_num} / {total}" if section else f"{slide_num} / {total}"
    _add_textbox(
        slide,
        MARGIN_X, FOOTER_TOP, SLIDE_W - MARGIN_X * 2, FOOTER_HEIGHT,
        text=text, font="Calibri", size=10, color=MUTED, align=PP_ALIGN.RIGHT,
    )


def _add_notes(slide, notes: str):
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


# --- Slide-type helpers -----------------------------------------------------


def add_title_slide(prs, title: str, subtitle: str, *, sub2: str = "", notes: str = "") -> None:
    slide = _blank_slide(prs)
    _set_background(slide, NAVY)
    _add_textbox(
        slide,
        MARGIN_X, Inches(2.0), SLIDE_W - MARGIN_X * 2, Inches(1.6),
        text=title, font="Calibri Light", size=56, bold=True,
        color=OFFWHITE, anchor=MSO_ANCHOR.MIDDLE,
    )
    _add_textbox(
        slide,
        MARGIN_X, Inches(3.7), SLIDE_W - MARGIN_X * 2, Inches(0.8),
        text=subtitle, font="Calibri", size=28, color=ACCENT,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    if sub2:
        _add_textbox(
            slide,
            MARGIN_X, Inches(4.6), SLIDE_W - MARGIN_X * 2, Inches(0.6),
            text=sub2, font="Calibri", size=18, color=OFFWHITE,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    _add_textbox(
        slide,
        MARGIN_X, Inches(6.6), SLIDE_W - MARGIN_X * 2, Inches(0.5),
        text=f"{TALK_VENUE}   ·   {TALK_DATE}   ·   {GITHUB_URL}",
        font="Calibri", size=12, color=MUTED,
    )
    _add_notes(slide, notes)


def add_bullets_slide(prs, title: str, bullets: list, *, notes: str = "",
                      section: str = "", slide_num: int = 0, total: int = 0) -> None:
    slide = _blank_slide(prs)
    _set_background(slide, OFFWHITE)
    _add_title(slide, title)
    box = slide.shapes.add_textbox(MARGIN_X, BODY_TOP, BODY_WIDTH, BODY_HEIGHT)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = level
        p.space_after = Pt(8)
        marker = "•  " if level == 0 else "–  "
        indent = "    " * level
        run = p.add_run()
        run.text = f"{indent}{marker}{text}" if text else ""
        run.font.name = "Calibri"
        run.font.size = Pt(22 if level == 0 else 18)
        run.font.color.rgb = DARK_TEXT if level == 0 else MUTED
    _add_footer(slide, slide_num, total, section)
    _add_notes(slide, notes)


def add_code_slide(prs, title: str, code: str, *, caption: str = "", notes: str = "",
                   section: str = "", slide_num: int = 0, total: int = 0,
                   language_hint: str = "csharp") -> None:
    slide = _blank_slide(prs)
    _set_background(slide, OFFWHITE)
    _add_title(slide, title)
    panel_h = BODY_HEIGHT - (Inches(0.4) if caption else Inches(0))
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_X, BODY_TOP, BODY_WIDTH, panel_h)
    panel.line.fill.background()
    panel.fill.solid()
    panel.fill.fore_color.rgb = CODE_BG
    panel.shadow.inherit = False
    tf = panel.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    lines = code.rstrip("\n").split("\n")
    if len(lines) <= 18:
        code_size = 16
    elif len(lines) <= 24:
        code_size = 14
    else:
        code_size = 12
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = "Consolas"
        run.font.size = Pt(code_size)
        run.font.color.rgb = CODE_TEXT
    if caption:
        _add_textbox(
            slide,
            MARGIN_X, Inches(6.65), BODY_WIDTH, Inches(0.4),
            text=caption, font="Calibri", size=14, color=MUTED,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    _add_footer(slide, slide_num, total, section)
    _add_notes(slide, notes)


def add_demo_slide(prs, title: str, steps: list, *, notes: str = "",
                   section: str = "", slide_num: int = 0, total: int = 0) -> None:
    slide = _blank_slide(prs)
    _set_background(slide, DEEP_BLUE)
    _add_textbox(
        slide,
        MARGIN_X, Inches(0.5), BODY_WIDTH, Inches(1.0),
        text="LIVE DEMO", font="Calibri Light", size=18, bold=True,
        color=ACCENT, anchor=MSO_ANCHOR.MIDDLE,
    )
    _add_textbox(
        slide,
        MARGIN_X, Inches(1.4), BODY_WIDTH, Inches(1.2),
        text=title, font="Calibri Light", size=44, bold=True,
        color=OFFWHITE, anchor=MSO_ANCHOR.MIDDLE,
    )
    box = slide.shapes.add_textbox(MARGIN_X, Inches(3.2), BODY_WIDTH, Inches(3.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, step in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = f"{i + 1}.  {step}"
        run.font.name = "Calibri"
        run.font.size = Pt(20)
        run.font.color.rgb = OFFWHITE
    _add_footer(slide, slide_num, total, section)
    notes_body = "In-game steps:\n" + "\n".join(f"  {i + 1}. {s}" for i, s in enumerate(steps))
    if notes:
        notes_body = notes + "\n\n" + notes_body
    _add_notes(slide, notes_body)


def add_quote_slide(prs, quote: str, attribution: str = "", *, notes: str = "",
                    section: str = "", slide_num: int = 0, total: int = 0) -> None:
    slide = _blank_slide(prs)
    _set_background(slide, NAVY)
    _add_textbox(
        slide,
        MARGIN_X, Inches(2.0), BODY_WIDTH, Inches(2.5),
        text=f"“{quote}”",
        font="Calibri Light", size=40,
        color=OFFWHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    if attribution:
        _add_textbox(
            slide,
            MARGIN_X, Inches(4.7), BODY_WIDTH, Inches(0.8),
            text=f"— {attribution}", font="Calibri", size=22, color=ACCENT,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
    _add_footer(slide, slide_num, total, section)
    _add_notes(slide, notes)


def add_compare_slide(prs, title: str, left_header: str, left_bullets: list,
                      right_header: str, right_bullets: list, *,
                      notes: str = "", section: str = "",
                      slide_num: int = 0, total: int = 0) -> None:
    """Two-column comparison. Used for the Kallikor / Timberborn intro."""
    slide = _blank_slide(prs)
    _set_background(slide, OFFWHITE)
    _add_title(slide, title)
    col_w = (BODY_WIDTH - Inches(0.4)) / 2
    col_h = BODY_HEIGHT - Inches(0.3)
    left_x = MARGIN_X
    right_x = MARGIN_X + col_w + Inches(0.4)
    for x, header, bullets, bg in (
        (left_x, left_header, left_bullets, COMPARE_LEFT_BG),
        (right_x, right_header, right_bullets, COMPARE_RIGHT_BG),
    ):
        panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, BODY_TOP, col_w, col_h)
        panel.line.fill.background()
        panel.fill.solid()
        panel.fill.fore_color.rgb = bg
        panel.shadow.inherit = False
        # Header bar.
        _add_textbox(
            slide, x + Inches(0.2), BODY_TOP + Inches(0.15), col_w - Inches(0.4), Inches(0.6),
            text=header, font="Calibri", size=22, bold=True, color=NAVY,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Bullets.
        box = slide.shapes.add_textbox(
            x + Inches(0.2), BODY_TOP + Inches(0.85),
            col_w - Inches(0.4), col_h - Inches(1.0),
        )
        tf = box.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(8)
            run = p.add_run()
            run.text = f"•  {b}"
            run.font.name = "Calibri"
            run.font.size = Pt(18)
            run.font.color.rgb = DARK_TEXT
    _add_footer(slide, slide_num, total, section)
    _add_notes(slide, notes)


# --- Sections + driver ------------------------------------------------------

HERE = Path(__file__).parent
SCREENSHOTS = HERE / "screenshots"

SECTIONS = {
    "open":     "Opening",
    "game":     "About Timberborn",
    "mod1":     "Mod #1: Quadruple Platform",
    "mod2a":    "Mod #2: Flood Season — flow",
    "mod2b":    "Mod #2: Flood Season — hazard",
    "cascade":  "The cascade",
    "save":     "Save / load",
    "mod3":     "Mod #3: Mixed Tide",
    "lessons":  "Gotchas & workflow",
    "mod4":     "Mod #4: Rainy Season (future)",
    "close":    "Closing",
}


@dataclass
class SlideContext:
    prs: Presentation
    section: str = ""
    slide_num: int = 0
    total: int = 0


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_builders = _slide_builders()
    total = len(slide_builders)
    ctx = SlideContext(prs=prs, total=total)
    for i, (section_key, fn) in enumerate(slide_builders, start=1):
        ctx.section = SECTIONS.get(section_key, "")
        ctx.slide_num = i
        fn(ctx)
    return prs


def _slide_builders():
    return [
        # Opening (4 slides ~ 5 min)
        ("open",     _slide_title),
        ("open",     _slide_kallikor_compare),
        ("open",     _slide_jasper_intro),
        ("open",     _slide_three_requests),

        # About Timberborn (4 slides ~ 6 min)
        ("game",     _slide_timberborn_60s),
        ("game",     _slide_demo_vanilla),
        ("game",     _slide_mod_loader),
        ("game",     _slide_three_layers),

        # Mod #1: Quadruple Platform (5 slides ~ 6 min)
        ("mod1",     _slide_mod1_intro),
        ("mod1",     _slide_mod1_templates),
        ("mod1",     _slide_mod1_blueprint),
        ("mod1",     _slide_mod1_timbermesh),
        ("mod1",     _slide_mod1_demo),

        # Mod #2 (a): Flow rate — clean seam (6 slides ~ 7 min)
        ("mod2a",    _slide_mod2a_intro),
        ("mod2a",    _slide_decomp_tool),
        ("mod2a",    _slide_iwaterstrength),
        ("mod2a",    _slide_modifier_class),
        ("mod2a",    _slide_bindito_wiring),
        ("mod2a",    _slide_demo_flow),

        # Mod #2 (b): Flood hazard — Harmony (5 slides ~ 7 min)
        ("mod2b",    _slide_mod2b_intro),
        ("mod2b",    _slide_randomizer_decomp),
        ("mod2b",    _slide_harmony_101),
        ("mod2b",    _slide_4line_postfix),
        ("mod2b",    _slide_demo_always_flood),

        # The cascade (4 slides ~ 6 min)
        ("cascade",  _slide_cascade_intro),
        ("cascade",  _slide_ui_helper_throws),
        ("cascade",  _slide_id_spoof),
        ("cascade",  _slide_inline_style_overrides),

        # Save / load (2 slides ~ 4 min)
        ("save",     _slide_save_and_postload),
        ("save",     _slide_fake_event_knockon),

        # Mod #3: Mixed Tide (4 slides ~ 6 min)
        ("mod3",     _slide_mod3_intro),
        ("mod3",     _slide_contamination_seam),
        ("mod3",     _slide_demo_mixed_tide),
        ("mod3",     _slide_doing_it_twice),

        # Gotchas catalogue (1 slide ~ 2 min)
        ("lessons",  _slide_gotcha_catalogue),

        # Mod #4 future (2 slides ~ 3 min)
        ("mod4",     _slide_mod4_rainy_intro),
        ("mod4",     _slide_mod4_design),

        # Close (1 slide + Q&A)
        ("close",    _slide_thanks_qa),
    ]


# --- Individual slides ------------------------------------------------------


def _slide_title(ctx: SlideContext) -> None:
    add_title_slide(
        ctx.prs,
        title="Modding Timberborn",
        subtitle=f"Three (and a half) mods I built with my {SON_AGE}-year-old",
        sub2=f"{PRESENTER} · designed by {SON_NAME} (age {SON_AGE})",
        notes=(
            "Opening line: 'This is a talk about modding Timberborn — a colony sim — "
            "with my 6-year-old as the customer. He couldn't read most of the code, but "
            "he made every design decision. Three shipped mods, one we haven't built "
            "yet, and a lot of gotchas in between.'\n\n"
            "Set expectations: technical walkthrough, three live demos, a future-work "
            "tease at the end, ~55 minutes content with Q&A after."
        ),
    )


def _slide_kallikor_compare(ctx: SlideContext) -> None:
    add_compare_slide(
        ctx.prs,
        title="Two simulators, side by side",
        left_header="What we build at Kallikor",
        left_bullets=[
            "Discrete-event simulation (jump from event to event, no fixed tick)",
            "Visualisation in PlayCanvas, decoupled from the sim (reads state, doesn't drive it)",
            "TypeScript / Python",
            "Sim and renderer can run on separate machines",
            "Domain: customer operations (logistics, warehouses, …)",
        ],
        right_header="What Timberborn is",
        right_bullets=[
            "Time-incrementing simulation (fixed tick rate drives state forward)",
            "Rendering in Unity, IN-PROCESS with the sim — every tick advances both",
            "C#  /  Unity 6 Mono runtime",
            "Sim and renderer share an object graph (Bindito DI)",
            "Domain: a beaver colony surviving seasonal hazards",
        ],
        notes=(
            "Walk this slide slowly — it's the conceptual anchor for the whole talk.\n\n"
            "Key things to highlight verbally:\n"
            "  - Discrete-event vs time-incrementing is a fundamental modelling choice. "
            "Time-incrementing is simpler to render (every frame IS a sim tick). "
            "Discrete-event lets you jump over uneventful time.\n"
            "  - PlayCanvas decoupling is a deliberate architectural choice we made — "
            "Timberborn made the opposite choice. Both are valid.\n"
            "  - C# .NET vs TypeScript: very different tooling stacks. Decompiling .NET "
            "is mature and easy — IL is well-specified and a tool like ilspycmd produces "
            "near-source-quality output. TypeScript bundles are messier to reverse-engineer.\n\n"
            "Lands the question: 'how do you change behaviour in a product whose stack "
            "is nothing like yours?'"
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_jasper_intro(ctx: SlideContext) -> None:
    add_bullets_slide(
        ctx.prs,
        title=f"The customer: {SON_NAME}, age {SON_AGE}",
        bullets=[
            f"{SON_NAME} loves Timberborn — has played it since he was four.",
            "Years before I had kids, I built Code for Life — a Blockly-based game that teaches kids to program.",
            (f"{SON_NAME} learned Blockly on it. Has started Python. Hadn't seen C# before this project.", 1),
            "We paired: I drove the keyboard and Claude Code. He drove the design.",
            ("Showed him Claude Code working alongside me — explained the agent's moves in plain English.", 1),
            ("Walked him through the actual C# at each step, in words he could follow.", 1),
            "His specs: behaviour-driven and short. 'Make floods big.'",
            "His feedback: fast and very honest. 'I don't like that.' / 'Cool, do more.'",
            "Best PM I've ever worked with.",
        ],
        notes=(
            "The human anchor of the talk. A few beats worth landing verbally:\n\n"
            "1. The Code for Life arc is genuinely worth a sentence. Built it before "
            "having kids → my kid grew up using it → he's grown enough to drive a real "
            "modding project. This talk is a continuation of that arc.\n\n"
            "2. The 'pair programming' framing matters: the actual coding pattern was "
            "I-on-keyboard-with-Claude-Code, Jasper next to me as customer + first "
            "user-tester. He didn't read most of the C# — he understood it at the "
            "level of 'this part picks the weather; this part draws the picture'.\n\n"
            "3. There's an obvious parallel between pairing with a 6-year-old and "
            "pairing with an AI agent: both are 'I'm in the driver's seat, you check "
            "my work and tell me what's next.' Worth flagging if the room is into it, "
            "but don't belabour. The talk is technical.\n\n"
            "4. Don't oversell the kid angle. He's not on stage. He shows up at each "
            "phase as a named customer."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_three_requests(ctx: SlideContext) -> None:
    add_bullets_slide(
        ctx.prs,
        title="Five asks — four answered, one pending",
        bullets=[
            f"#1  '{JASPER_PLATFORM_BRIEF}'  → Quadruple Platform (templates)",
            f"#2  '{JASPER_TEMPERATE_QUOTE}'  → Flood Season A: more flow in temperate (interface)",
            f"#3  '{JASPER_HAZARD_QUOTE}'  → Flood Season B: flood as a new hazard (Harmony)",
            f"#4  '{JASPER_MIXED_QUOTE}'  → Mixed Tide (Harmony, parallel branches)",
            f"#5  '{JASPER_RAIN_BRIEF}'  → Rainy Season (not built yet)",
            "",
            ("Each ask escalates the modding machinery. We'll walk through layer by layer.", 1),
        ],
        notes=(
            "Set up the talk's structure here. The audience should leave this slide "
            "with the whole running order in their head:\n\n"
            "  #1  Quadruple Platform — pure data, no compilation. Layer 1: templates.\n"
            "  #2  Flood Season A — implement an exposed interface, register a settings UI. Layer 2.\n"
            "  #3  Flood Season B — replace a hardcoded vanilla return value with a new type. Layer 3: Harmony.\n"
            "  #4  Mixed Tide — same patches as Flood Season, parallel branches. Still layer 3, stress-test.\n"
            "  #5  Rainy Season — future-work tease.\n\n"
            "Note the two halves of Flood Season — requests #2 and #3 landed in the same "
            "mod but as two distinct phases driven by two distinct asks from Jasper.\n\n"
            "Asks escalate roughly: data change → interface → patch → patch-cascade → unknown."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_timberborn_60s(ctx: SlideContext) -> None:
    add_bullets_slide(
        ctx.prs,
        title="Timberborn in 60 seconds",
        bullets=[
            "Beaver colony sim by Mechanistry. Steam early access since 2021.",
            "Build dams, store water, plant trees, survive seasonal hazards.",
            "Three weather phases per cycle: Temperate → Drought OR Badtide (contaminated water).",
            "Built on Unity 6, written in C#. Time-incrementing sim drives the render every tick.",
            ("No Unity Editor needed for modding — mods load as compiled DLLs into the shipping runtime.", 1),
        ],
        notes=(
            "~45 seconds here. Audience needs to know:\n"
            "  - There are WATER SOURCES, BEAVERS, and HAZARDS\n"
            "  - The game has weather cycles\n"
            "  - It's Unity + C#\n"
            "Move on quickly — the rest of the talk is the modding."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_demo_vanilla(ctx: SlideContext) -> None:
    add_demo_slide(
        ctx.prs,
        title="This is Timberborn",
        steps=[
            "Switch to the game — paused at a save with an active Folktails base.",
            "Point at: water sources flowing, beavers working, the top-right date panel.",
            "Pan to a body of water — the column simulation is doing its thing.",
            "If a hazard is approaching, point at the forecast banner. Otherwise mention vanilla cycles drought / badtide / temperate.",
            "Switch back to slides — 'OK, now we mod it.'",
        ],
        notes=(
            "60-90 seconds. Don't narrate the gameplay in depth — just orient the "
            "audience. The bits that matter for the rest of the talk:\n"
            "  - Water sources are a thing\n"
            "  - The date panel + weather panel sit top-right (we mod both)\n"
            "  - Cycles tick through Temperate → Hazard → Temperate\n\n"
            "Don't get drawn into beaver lore. Move on."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_mod_loader(ctx: SlideContext) -> None:
    add_code_slide(
        ctx.prs,
        title="How Timberborn loads mods",
        code=(
            "Documents/Timberborn/Mods/Flood Season/\n"
            "    manifest.json     # name, version, required-mods\n"
            "    Code.dll          # compiled C# (optional — some mods are data-only)\n"
            "    assets/           # PNGs, blueprints, .timbermesh, .csv localisations\n"
            "\n"
            "// On launch:\n"
            "var dlls = modDir.GetFiles(\"*.dll\", SearchOption.AllDirectories);  // RECURSIVE\n"
            "foreach (var dll in dlls) {\n"
            "    var asm = Assembly.LoadFile(dll.FullName);\n"
            "    foreach (var t in asm.GetTypes()) {\n"
            "        if (typeof(IModStarter).IsAssignableFrom(t)) {\n"
            "            ((IModStarter)Activator.CreateInstance(t)).StartMod(env);\n"
            "        }\n"
            "    }\n"
            "}\n"
            "// JSON blueprints in /assets/ are discovered separately via the templates system."
        ),
        caption="No installer. The mod IS the folder. Recursive scan is the source of a future gotcha.",
        notes=(
            "Two ideas to land:\n"
            "1. Reflection-based discovery. No registration manifest, no init order. The "
            "game just walks the folder and instantiates anything that looks like a mod.\n"
            "2. 'RECURSIVE' is doing work in that sentence — we'll come back to this when "
            "bin/obj/Code.dll bites us during the build.\n\n"
            "Some mods are PURE DATA (no Code.dll). The Quadruple Platform we're about to "
            "show is one of those."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_three_layers(ctx: SlideContext) -> None:
    add_bullets_slide(
        ctx.prs,
        title="Three layers of modding (in increasing difficulty)",
        bullets=[
            "1. Templates — declarative JSON blueprints. No code. Game's content schema.",
            ("Add buildings / change costs / append to lists. Mod #1: Quadruple Platform.", 1),
            "2. Exposed interfaces — plug a class in via DI. The game asks you for behaviour.",
            ("Game has thought about extension here. Mod #2 part A: flow modifier.", 1),
            "3. Harmony patches — rewrite shipped methods at runtime.",
            ("No seam exists. You're the uninvited guest. Mod #2 part B + Mod #3.", 1),
            "",
            ("Reach for layer N only when layer N−1 can't do the job.", 1),
        ],
        notes=(
            "The conceptual scaffolding for the rest of the talk. Each mod we walk "
            "through sits at one or more of these layers. The pedagogical order — "
            "templates first, Harmony last — matches the order of the customer "
            "requests, which is a nice accident.\n\n"
            "Land the rule of thumb: 'reach for layer N only when layer N−1 can't do the "
            "job.' It's not just modding hygiene; it's how to think about any "
            "third-party extension."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_mod1_intro(ctx: SlideContext) -> None:
    add_bullets_slide(
        ctx.prs,
        title=f"Mod #1: Quadruple Platform — '{JASPER_PLATFORM_BRIEF}'",
        bullets=[
            "Timberborn ships Single, Double, and Triple platforms. Stack-three is the limit.",
            f"{SON_NAME} wanted to build a FOUR-high platform.",
            "First instinct: 'this is a data change, not a code change.'",
            ("The game's content (buildings, recipes, costs, beaver needs) lives in JSON blueprints.", 1),
            ("If we can append a new blueprint and a new 3D model, we're done.", 1),
            "Lives at Mods/Quadruple Platform/ — entirely separate from the Flood Season mod.",
        ],
        notes=(
            "Set up the simplest mod first. The audience should leave this section "
            "thinking 'oh, lots of game content is just JSON — modding doesn't have to "
            "mean Harmony.'\n\n"
            "Mention: Timberborn's content pipeline is a deliberate design choice by "
            "Mechanistry. They want modders. Not every game does this — some games ship "
            "everything packed into binary asset bundles."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_mod1_templates(ctx: SlideContext) -> None:
    add_code_slide(
        ctx.prs,
        title="Templates: how the game knows about buildings",
        code=(
            "// TemplateCollections/TemplateCollection.Buildings.Folktails.blueprint.json\n"
            "{\n"
            "  \"TemplateCollectionSpec\": {\n"
            "    \"CollectionId\": \"Buildings.Folktails\",\n"
            "    \"Blueprints#append\": [                  // ← the magic suffix\n"
            "      \"Buildings/Paths/QuadruplePlatform/\n"
            "       QuadruplePlatform.Folktails.blueprint\"\n"
            "    ]\n"
            "  }\n"
            "}\n"
            "\n"
            "// '#append' tells the game's content loader: don't OVERRIDE the\n"
            "// vanilla list, EXTEND it. We're adding to the Folktails buildings\n"
            "// catalogue without touching anything else."
        ),
        caption="The `#append` suffix is the whole story. Declarative list-extension via JSON.",
        notes=(
            "This is the conceptual win for the templates layer. The game's content "
            "pipeline lets us EXTEND vanilla lists without forking them. No code, no "
            "Harmony — just a key suffix.\n\n"
            "Worth mentioning: there are other suffixes too (`#replace`, `#remove`) but "
            "`#append` is the one we needed. The templates system is fully documented "
            "by Mechanistry."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_mod1_blueprint(ctx: SlideContext) -> None:
    add_code_slide(
        ctx.prs,
        title="The Quadruple Platform blueprint",
        code=(
            "// Buildings/.../QuadruplePlatform.Folktails.blueprint.json (excerpt)\n"
            "\"BuildingSpec\": {\n"
            "    \"BuildingCost\": [{\"Id\": \"Plank\", \"Amount\": 12}],\n"
            "    \"ScienceCost\": 250,\n"
            "    …\n"
            "},\n"
            "\"BlockObjectSpec\": {\n"
            "    \"Size\": {\"X\": 1, \"Y\": 1, \"Z\": 4},     // 1x1 footprint, 4 high\n"
            "    \"Blocks\": [ /* 4 stacked block descriptors */ ],\n"
            "    …\n"
            "},\n"
            "\"LabeledEntitySpec\": {\n"
            "    \"DisplayNameLocKey\": \"Building.QuadruplePlatform.DisplayName\",\n"
            "    \"Icon\": \"Buildings/.../QuadruplePlatformIcon\"\n"
            "},\n"
            "\"Children\": {\n"
            "    \"#Finished\":         { \"TimbermeshSpec\": { \"Model\": \"…\" } },\n"
            "    \"#FinishedUncovered\": { \"TimbermeshSpec\": { \"Model\": \"…\" } },\n"
            "    \"#Unfinished\":        { /* construction stages */ }\n"
            "}"
        ),
        caption="Cost, footprint, name, icon, models. 200-ish lines of JSON. No code.",
        notes=(
            "Audience reaction here should be 'oh, that's basically a spec sheet'. The "
            "whole building — its cost, its footprint, its models, its localisation key "
            "— is data.\n\n"
            "Worth pointing out: the blueprint references three separate 3D models — "
            "Finished, FinishedUncovered (no roof), and Unfinished (construction site). "
            "All Folktails-themed."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_mod1_timbermesh(ctx: SlideContext) -> None:
    add_code_slide(
        ctx.prs,
        title="The 3D model: protobuf, zlib, no Blender",
        code=(
            "# Mods/Quadruple Platform/generate_models.py  (excerpt)\n"
            "# The .timbermesh format is zlib-compressed Protocol Buffers (proto3).\n"
            "# Mechanistry publishes the .proto schema on GitHub:\n"
            "#   https://github.com/mechanistry/timbermesh/blob/main/proto/model.proto\n"
            "\n"
            "def build_node(name, parent, position, rotation, scale,\n"
            "               vertex_count, vertex_properties, meshes):\n"
            "    result = b\"\"\n"
            "    result += field_varint(1, parent)\n"
            "    result += field_string(2, name)\n"
            "    result += field_message(3, build_vector3(*position))\n"
            "    result += field_message(4, build_quaternion(*rotation))\n"
            "    result += field_message(5, build_vector3(*scale))\n"
            "    result += field_varint(6, vertex_count)\n"
            "    for vp in vertex_properties: result += field_message(7, vp)\n"
            "    for mesh in meshes:          result += field_message(8, mesh)\n"
            "    return result\n"
            "\n"
            "# I described what I wanted: \"a Quadruple Platform, mirroring the\n"
            "# Triple, four pillars and a top slab.\" Claude wrote the whole\n"
            "# script — protobuf encoding included. It picked its own textures.\n"
            "# (Reader: those textures are... distinctive.)"
        ),
        caption="AI wrote the script from a description. AI also picked the textures. Result: funky.",
        notes=(
            "The reveal: the .timbermesh model is zlib+protobuf, schema published on "
            "GitHub. You don't need to model in Blender and export — you can build the "
            "binary directly from a script.\n\n"
            "Be honest about the AI involvement — I didn't write this Python. I "
            "described what I wanted (a Quadruple Platform mirroring the Triple), and "
            "Claude wrote the whole thing: protobuf encoding, geometry, materials, "
            "everything. It also picked its own textures from the materials it found in "
            "the base game's models. The result is a working four-high platform that "
            "looks… let's say 'distinctive'. Show the in-game screenshot on the demo "
            "slide if you want the laugh.\n\n"
            "Don't apologise for the AI-wrote-it angle. It's the point. This is what "
            "'AI-driven reverse-engineering of binary formats' actually looks like in "
            "practice — describe the goal, judge the output, iterate. Same workflow "
            "would apply to PlayCanvas asset bundles in our stack."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_mod1_demo(ctx: SlideContext) -> None:
    add_demo_slide(
        ctx.prs,
        title="This is with a Quadruple Platform",
        steps=[
            "Switch to Timberborn — Folktails base loaded, building menu open.",
            "Paths group → Quadruple Platform appears with its (AI-textured) icon.",
            "Place one on flat ground. Beavers haul planks. Construction stages through.",
            "Stand the platform next to a vanilla Triple — the 4-high difference is obvious.",
            "Pause on the textures for the laugh.",
        ],
        notes=(
            "The point of this demo is to convince the audience that a JSON-only mod "
            "is GAME-COMPLETE. There's no asterisk — building shows up in the toolbar, "
            "beavers build it, it works.\n\n"
            "Land the funky-textures payoff here. The AI picked colours from the base "
            "game's material palette and combined them in an order Mechanistry would "
            "not have. Worth a small laugh."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_mod2a_intro(ctx: SlideContext) -> None:
    add_quote_slide(
        ctx.prs,
        quote=JASPER_TEMPERATE_QUOTE,
        attribution=f"{SON_NAME}, feature request #2",
        notes=(
            "Land the quote, pause, then: 'OK, what does \"too much water\" mean? "
            "Let's start small — can we make the wet seasons WETTER? Double the flow "
            "during temperate? That's a one-line spec we can act on. Floods themselves "
            "come later, as a separate ask.'\n\n"
            "Important framing: this is the SCOPED interpretation of his request. He "
            "wasn't being specific about 'flood' yet — he just wanted more water. The "
            "temperate-multiplier story IS the response to this ask. Request #3 "
            "(flood-as-hazard) came LATER, separately."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_decomp_tool(ctx: SlideContext) -> None:
    add_code_slide(
        ctx.prs,
        title="Step zero: read the source you don't have",
        code=(
            "$ dotnet tool install -g ilspycmd\n"
            "$ ilspycmd -p -o ~/decomp/WaterSourceSystem \\\n"
            "    \"…/Timberborn_Data/Managed/Timberborn.WaterSourceSystem.dll\"\n"
            "\n"
            "# Result: a tree of readable .cs files. Mechanistry doesn't obfuscate.\n"
            "# Greppable. Diff-able. Not source-of-truth, but very close.\n"
            "\n"
            "$ grep -rn \"IWaterStrengthModifier\" decomp/WaterSourceSystem/\n"
            "    IWaterStrengthModifier.cs:5:    public interface IWaterStrengthModifier\n"
            "    WaterSource.cs:42:    private readonly List<IWaterStrengthModifier>\n"
            "    WaterSource.cs:149:    foreach (IWaterStrengthModifier waterStrengthModifier"
        ),
        caption="The whole workflow: decompile, grep, find the seam.",
        notes=(
            "This is the skill that transfers to all third-party software work. Install "
            "the decompiler, point it at the binary, get readable code, grep for what "
            "you care about.\n\n"
            ".NET is the easy case — IL is structured, ilspycmd is mature. For "
            "TypeScript bundles you'd reach for sourcemap inspection or a tool like "
            "Webcrack. Different tooling, same workflow."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_iwaterstrength(ctx: SlideContext) -> None:
    add_code_slide(
        ctx.prs,
        title="The seam: IWaterStrengthModifier",
        code=(
            "// Timberborn.WaterSourceSystem (vanilla, decompiled)\n"
            "public interface IWaterStrengthModifier {\n"
            "    float GetStrengthModifier();\n"
            "}\n"
            "\n"
            "// WaterSource ticks every game tick and composes ALL modifiers:\n"
            "private void UpdateCurrentStrength() {\n"
            "    float num = SpecifiedStrength;\n"
            "    foreach (IWaterStrengthModifier m in _waterStrengthModifiers) {\n"
            "        num *= m.GetStrengthModifier();   // multiplicative chain\n"
            "    }\n"
            "    SetCurrentStrength(num);\n"
            "}"
        ),
        caption="Drought ships its own IWaterStrengthModifier. We just register one more.",
        notes=(
            "Land the WIN. We have:\n"
            "  - a one-method interface\n"
            "  - a multiplicative composition chain (we don't fight other modifiers)\n"
            "  - a per-entity registration point (AddWaterStrengthModifier)\n\n"
            "This is what 'exposed interface' looks like. The game's authors wanted us "
            "to be able to do this. Drought is implemented the same way, internally."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_modifier_class(ctx: SlideContext) -> None:
    add_code_slide(
        ctx.prs,
        title="Our modifier — boost flow during temperate",
        code=(
            "internal class FloodSeasonWaterStrengthModifier\n"
            "    : BaseComponent, IAwakableComponent, IInitializableEntity,\n"
            "      IWaterStrengthModifier {\n"
            "\n"
            "    private readonly WeatherService _weatherService;\n"
            "    private readonly FloodSeasonSettings _settings;\n"
            "\n"
            "    public float GetStrengthModifier() {\n"
            "        if (_weatherService.IsHazardousWeather) {\n"
            "            return 1.0f;   // hazardous phase: don't touch flow\n"
            "        }\n"
            "        return _settings.TemperateMultiplier;   // wetter wet seasons\n"
            "    }\n"
            "}"
        ),
        caption="That's the whole mechanic. ~15 lines. We'll extend the class later as new requests land.",
        notes=(
            "Walk the code top to bottom. Things the audience should notice:\n"
            "  - Standard mod-component shape: BaseComponent + IWaterStrengthModifier\n"
            "    + the lifecycle interfaces Timberborn expects (Awake, InitializeEntity).\n"
            "  - The gate on _weatherService.IsHazardousWeather — during droughts and\n"
            "    badtides we return 1.0 and let the game's own modifiers handle it.\n"
            "  - The temperate path returns the user-configured multiplier from the\n"
            "    settings panel we're about to see on the next slide.\n\n"
            "This is the WHOLE first half of the Flood Season mod. We'll extend this\n"
            "class as new asks land (the switch with FloodWeather / MixedTideWeather\n"
            "cases shows up later, once those types exist) — but right now, this is it."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_bindito_wiring(ctx: SlideContext) -> None:
    add_code_slide(
        ctx.prs,
        title="Wiring it all up: Bindito DI + TemplateModule + ModSettings",
        code=(
            "[Context(\"Game\")]                       // auto-discovered by Bindito\n"
            "internal class FloodSeasonConfigurator : Configurator {\n"
            "    protected override void Configure() {\n"
            "        Bind<FloodSeasonSettings>().AsSingleton();    // ← UI panel\n"
            "        Bind<FloodSeasonWaterStrengthModifier>().AsTransient();\n"
            "        MultiBind<TemplateModule>()\n"
            "            .ToProvider(ProvideTemplateModule).AsSingleton();\n"
            "    }\n"
            "    private static TemplateModule ProvideTemplateModule() {\n"
            "        var b = new TemplateModule.Builder();\n"
            "        // \"For every WaterSource entity spawned, attach a modifier.\"\n"
            "        b.AddDecorator<WaterSource, FloodSeasonWaterStrengthModifier>();\n"
            "        return b.Build();\n"
            "    }\n"
            "}\n"
            "// FloodSeasonSettings : ModSettingsOwner (from eMka.ModSettings)\n"
            "//   public ModSetting<int> TemperateMultiplierPercent { get; } = ...;\n"
            "//   UI controls auto-rendered from public ModSetting<T> properties."
        ),
        caption="DI + decorator pattern + reflection-driven settings UI. Three layers in 15 lines.",
        notes=(
            "Three concepts crammed into one slide because they're each one paragraph:\n"
            "  1. Bindito DI — [Context(\"Game\")] auto-discovery, AsSingleton / AsTransient.\n"
            "  2. TemplateModule — decorator wiring 'attach component T2 wherever T1 exists'.\n"
            "  3. eMka.ModSettings — third-party mod that reflects over ModSetting<T> properties to generate UI controls.\n\n"
            "If anyone asks 'what's eMka.ModSettings?' — it's a Workshop mod that lots "
            "of Timberborn mods depend on for a consistent settings UI. We declare it "
            "as required in manifest.json."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_demo_flow(ctx: SlideContext) -> None:
    add_demo_slide(
        ctx.prs,
        title="This is with temperate flow increased",
        steps=[
            "Open Mod Settings panel in-game.",
            "Set 'Temperate flow multiplier (%)' to 500.",
            "Close settings. Point at the water sources — flow visibly jumps.",
            "Set back to 200 (default 2×) to keep things sane.",
            "Optional: speed-time to make the downstream effect obvious.",
        ],
        notes=(
            "Pre-load a save with several water sources visible — an empty map makes "
            "for a boring demo.\n\n"
            "Talking point during demo: 'This is the WHOLE first half of Flood Season "
            "in 15 lines of C# plus a settings property. We didn't need Harmony at all. "
            "The game gave us the seam.'"
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_mod2b_intro(ctx: SlideContext) -> None:
    add_quote_slide(
        ctx.prs,
        quote=JASPER_HAZARD_QUOTE,
        attribution=f"{SON_NAME}, feature request #3",
        notes=(
            "The escalation. Wetter wet seasons (Flood Season A) was great — and then "
            "Jasper came back specifically asking for a flood SEASON, like the droughts "
            "and badtides he already knew about. Not just more water — a whole hazard.\n\n"
            "Speaker bridge to the next slide:\n"
            "  'Translation: introduce a NEW IHazardousWeather type, alongside Drought "
            "and Badtide. Look for a seam… and there isn't one. The picker is a "
            "hardcoded if/else with two private fields. No MultiBind, no event hook, "
            "no registration call. Welcome to layer 3 — Harmony.'\n\n"
            "Then move straight to the next slide which shows the decompiled if/else.\n\n"
            "Tone shift the audience should feel: up to here, modding has felt "
            "'supported'. From here on, we're rewriting compiled methods at runtime."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_randomizer_decomp(ctx: SlideContext) -> None:
    add_code_slide(
        ctx.prs,
        title="HazardousWeatherRandomizer — no seam to hook",
        code=(
            "// Timberborn.HazardousWeatherSystem (vanilla, decompiled)\n"
            "public IHazardousWeather GetRandomWeatherForCycle(int cycle) {\n"
            "    if (ShouldBeBadtideWeather(cycle)) {\n"
            "        return _badtideWeather;\n"
            "    }\n"
            "    return _droughtWeather;\n"
            "}\n"
            "\n"
            "// That's it. No interface enumeration. No event hook.\n"
            "// Just an if/else with two hardcoded private fields."
        ),
        caption="The thing we want to extend is the most closed shape in the whole game.",
        notes=(
            "Land this hard. The audience should appreciate that the game's authors "
            "didn't WANT this to be extensible. They wrote drought-or-badtide and moved "
            "on.\n\n"
            "We're about to overwrite the return value of this method at runtime, with "
            "no permission. Welcome to Harmony."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_harmony_101(ctx: SlideContext) -> None:
    add_bullets_slide(
        ctx.prs,
        title="Harmony — the runtime patcher",
        bullets=[
            "Library that rewrites method IL at load time. Originally for RimWorld mods.",
            "You don't ship modified game DLLs. You ship patch INSTRUCTIONS.",
            "Three patch positions:",
            ("PREFIX — runs before the original. Can read args. `return false` SKIPS the original.", 1),
            ("POSTFIX — runs after. Can read AND modify the original's return value via `ref __result`.", 1),
            ("TRANSPILER — rewrites the original's IL itself. We didn't need this.", 1),
            "Patches discovered by [HarmonyPatch] attribute + a single PatchAll() call at startup.",
        ],
        notes=(
            "If anyone in the audience knows Harmony from RimWorld, Stardew, BepInEx, "
            "Cities Skylines, Subnautica — same library, same mental model. 'Don't "
            "modify the binary; ship a recipe for patching it at load.'\n\n"
            "Mention: Harmony isn't bundled with Timberborn. Our mod's manifest "
            "declares 'Harmony for Timberborn' as a required dependency, which "
            "redistributes 0Harmony.dll legally."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_4line_postfix(ctx: SlideContext) -> None:
    add_code_slide(
        ctx.prs,
        title="The hijack — a Harmony postfix in 10 lines",
        code=(
            "[HarmonyPatch(typeof(HazardousWeatherRandomizer),\n"
            "              nameof(HazardousWeatherRandomizer.GetRandomWeatherForCycle))]\n"
            "internal static class RandomizerPatch {\n"
            "\n"
            "    [HarmonyPostfix]\n"
            "    public static void Postfix(int cycle, ref IHazardousWeather __result) {\n"
            "        // … gates omitted: enable, grace cycles, weighted roll …\n"
            "        if (RollSaysFlood()) {\n"
            "            __result = FloodWeather.Instance;\n"
            "        } else if (RollSaysMixedTide()) {\n"
            "            __result = MixedTideWeather.Instance;\n"
            "        }\n"
            "    }\n"
            "}"
        ),
        caption="Vanilla picks drought / badtide. We OVERWRITE the result. Done.",
        notes=(
            "Land the surprise: this is all it takes to extend a method that has no "
            "extension point. The original runs (so the game's internal streak tracking "
            "still ticks normally), then we get a chance to overwrite the return value.\n\n"
            "Quick note on `__result`: that's a Harmony magic parameter. Two underscores. "
            "The original method's return slot, available to read and write. Harmony has "
            "other magic parameters too — `__instance`, `___fieldName` — we'll see one of "
            "those in a couple of slides.\n\n"
            "FloodWeather.Instance: Harmony patches are static methods, so they can't "
            "take DI. We publish a static accessor from the Bindito-managed FloodWeather "
            "constructor. Pattern: stash the singleton on a static field."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_demo_always_flood(ctx: SlideContext) -> None:
    add_demo_slide(
        ctx.prs,
        title="This is with a flood season",
        steps=[
            "Open Mod Settings, enable Flood Season, probability 100%, grace 0.",
            "Start a NEW game — cycle weather is decided at cycle START (gotcha).",
            "Fast-forward to the first hazardous cycle.",
            "Point at the panel: 'A flood has begun!' instead of drought/badtide.",
            "Note that the icon + banner still look like drought (cascade fixes coming).",
        ],
        notes=(
            "Don't try to demo on an existing save — the cycle's weather was already "
            "decided as something else. Settings changes never retro-apply mid-cycle. "
            "New game.\n\n"
            "Talking point during the demo: 'Notice the icon, the banner, the labels — "
            "they all still say drought/badtide. The game LOGIC thinks it's a flood. "
            "The UI hasn't caught up. That's the next twenty minutes of the talk.'"
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_cascade_intro(ctx: SlideContext) -> None:
    add_bullets_slide(
        ctx.prs,
        title="And then everything else breaks",
        bullets=[
            "We introduced a NEW IHazardousWeather type. Vanilla code didn't expect this.",
            "Every place that does `is DroughtWeather else is BadtideWeather else throw` fires.",
            "Every place keyed off the weather's string Id misses the lookup (fog spec, history).",
            "Every UI element CSS-bound to drought/badtide classes shows the wrong art.",
            "Each is a separate patch. We discovered them by crashing.",
            ("Walking through the categories with one representative fix each.", 1),
        ],
        notes=(
            "The slide that motivates the rest of the section. The audience should walk "
            "away with: 'introducing one new type into a hardcoded domain causes a "
            "cascade. The cost of bypassing the seam isn't four lines, it's the long "
            "tail of vanilla code that fights back.'\n\n"
            "We'll show the throw-on-unknown-type pattern (UI helper), the Id-spoof trick "
            "for string-keyed lookups, then the inline-style-override pattern for the UI."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_ui_helper_throws(ctx: SlideContext) -> None:
    add_code_slide(
        ctx.prs,
        title="The first throw: HazardousWeatherUIHelper",
        code=(
            "// Vanilla — Timberborn.HazardousWeatherSystemUI (decompiled)\n"
            "private void UpdateCurrentUISpecification() {\n"
            "    var current = _hazardousWeatherService.CurrentCycleHazardousWeather;\n"
            "    if (!(current is DroughtWeather)) {\n"
            "        if (!(current is BadtideWeather)) {\n"
            "            throw new InvalidOperationException(\n"
            "                \"No UI for weather: \" + current);\n"
            "        }\n"
            "        _currentUISpecification = _badtideWeatherUISpecification;\n"
            "    } else {\n"
            "        _currentUISpecification = _droughtWeatherUISpecification;\n"
            "    }\n"
            "}\n"
            "// Fix: Harmony PREFIX that substitutes the drought UI spec when our flood\n"
            "// is current, then returns false to skip the original method entirely."
        ),
        caption="Three throws like this in vanilla. UI helper, sound player, fog spec dictionary.",
        notes=(
            "The fix is a Harmony PREFIX that returns false (skipping the original) "
            "when our flood is current, substituting drought's UI spec in place of the "
            "would-have-thrown branch.\n\n"
            "Sound player and fog spec have the same shape — different fix specifics. "
            "The fog spec is keyed by string Id, which leads to the next slide: the "
            "Id-spoof trick."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_id_spoof(ctx: SlideContext) -> None:
    add_code_slide(
        ctx.prs,
        title="The Id-spoof trick — lie on purpose",
        code=(
            "internal class FloodWeather : IHazardousWeather, ILoadableSingleton {\n"
            "\n"
            "    public string Id => \"DroughtWeather\";  // ← deliberate\n"
            "\n"
            "    // Vanilla code does:\n"
            "    //   var fog = Sun.GetFogSettings(currentWeather.Id);\n"
            "    //   var count = HazardousWeatherHistory.GetCyclesCount(id);\n"
            "    // Both are string-keyed dictionary lookups. Blueprint data ships only\n"
            "    // for \"DroughtWeather\" and \"BadtideWeather\". A novel id would throw\n"
            "    // KeyNotFoundException.\n"
            "    //\n"
            "    // Returning \"DroughtWeather\" → the lookup hits the drought entry.\n"
            "    // Fog thinks we're a drought. History counts us as a drought.\n"
            "    //\n"
            "    // Type-based dispatch (`is FloodWeather`) STILL works — that's C# type\n"
            "    // identity, not the string. Our own modifier still distinguishes us.\n"
            "    // …\n"
            "}\n"
            "// Mixed Tide does the same trick but spoofs as \"BadtideWeather\"."
        ),
        caption="Spoof the id where vanilla uses strings. Keep the type where vanilla uses types.",
        notes=(
            "Slow down for this slide — it's one of the cleverer bits.\n\n"
            "The trick:\n"
            "  - For string-keyed lookups (fog spec, history count), lie about the id. "
            "Inherit drought's data.\n"
            "  - For type-based dispatch (`is FloodWeather`), the C# type stays distinct. "
            "Our modifier still tells flood from drought via type identity.\n\n"
            "The side effect — vanilla's HazardousWeatherHistory counts our floods as "
            "droughts and very slightly affects its handicap math — is acceptable. "
            "Mention as a known wart."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_inline_style_overrides(ctx: SlideContext) -> None:
    add_bullets_slide(
        ctx.prs,
        title="UI: drought textures bleed through, we override inline",
        bullets=[
            "Spoofing the UI spec stops vanilla throwing — but the player sees DROUGHT art during a FLOOD.",
            "Date panel icon: CSS class on root → swap inline style.backgroundImage.",
            "Notification banner: same. Find the Image element, override inline.",
            "Weather panel progress strip: it's a CUSTOM MESH, not a regular CSS background.",
            ("Sprite is loaded via Resources.Load from a CustomStyleProperty<string>.", 1),
            ("UI Toolkit resolution is DEFERRED — our override gets undone after our postfix.", 1),
            ("Fix: postfix the resolver event itself, so we write _image LAST.", 1),
        ],
        notes=(
            "Three UI surfaces, three different fix strategies.\n\n"
            "The first two are 'inline override of a CSS-driven background-image' — "
            "straightforward UI Toolkit.\n\n"
            "The third one (the progress strip) is the deferred-resolver race. We "
            "discovered this by adding Debug.Log diagnostics to two candidate patches — "
            "one log fired but the visual didn't change, telling us the resolver was "
            "running AFTER our postfix and overwriting our work.\n\n"
            "Lesson: UI Toolkit's custom style resolution is queued, not synchronous. "
            "Race conditions are real even in single-threaded UI code."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_save_and_postload(ctx: SlideContext) -> None:
    add_code_slide(
        ctx.prs,
        title="Save / load: vanilla schema can't fit",
        code=(
            "// Vanilla persists the active hazard as ONE BOOL:\n"
            "//   IsDrought: true  (Drought)\n"
            "//   IsDrought: false (Badtide)\n"
            "// No room for a third state. Save during a flood → loads back as badtide.\n"
            "\n"
            "internal class HazardousWeatherStatePersistence\n"
            "    : ILoadableSingleton, ISaveableSingleton, IPostLoadableSingleton {\n"
            "\n"
            "    private static readonly SingletonKey StateKey =\n"
            "        new SingletonKey(\"Spycho.FloodSeason.HazardousWeatherState\");\n"
            "    private static readonly PropertyKey<bool> IsFloodActiveKey      = …;\n"
            "    private static readonly PropertyKey<bool> IsMixedTideActiveKey  = …;\n"
            "\n"
            "    public void Save(ISingletonSaver s) {\n"
            "        var saver = s.GetSingleton(StateKey);\n"
            "        saver.Set(IsFloodActiveKey,     current is FloodWeather);\n"
            "        saver.Set(IsMixedTideActiveKey, current is MixedTideWeather);\n"
            "    }\n"
            "    public void PostLoad() {\n"
            "        // Read bools → force-restore the right weather via AccessTools.PropertySetter\n"
            "        // → re-fire HazardousWeatherSelectedEvent (UI refresh)\n"
            "        // → post fake HazardousWeatherEndedEvent(badtide) to undo entity init.\n"
            "    }\n"
            "}"
        ),
        caption="Side-car schema, AccessTools setter override, synthetic event for entity cleanup.",
        notes=(
            "Three ideas crammed into one slide:\n"
            "1. Vanilla's bool can't represent our state. So we side-car: our own "
            "singleton-keyed save data alongside vanilla's.\n"
            "2. The setter for CurrentCycleHazardousWeather is private. AccessTools."
            "PropertySetter (a Harmony utility) bypasses the access check.\n"
            "3. After we force-set the right weather, we post a synthetic 'badtide "
            "ended' event so vanilla's entity-side contamination controllers (which "
            "vanilla load enabled by mistake) walk through their disable handler.\n\n"
            "This sets up the next slide — the fake-event knock-on."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_fake_event_knockon(ctx: SlideContext) -> None:
    add_bullets_slide(
        ctx.prs,
        title="The fake-event knock-on",
        bullets=[
            "Posting HazardousWeatherEndedEvent(BadtideWeather) cleans up contamination ✓",
            "It ALSO reaches GameMusicPlayer.OnHazardousWeatherEnded ✗",
            ("Music player thinks badtide just ended → stops drought track, starts standard.", 1),
            ("But Load() already started drought music. Standard plays ON TOP of drought.", 1),
            "Fix: Harmony PREFIX on GameMusicPlayer.OnHazardousWeatherEnded.",
            ("If ended-weather is badtide AND our custom hazard is current → skip handler.", 1),
            ("That combination is unreachable in vanilla. Only our synthetic post produces it.", 1),
            "Lesson: synthetic events leak across every subscriber on the bus. No private channel.",
        ],
        notes=(
            "Audience reaction here should be a wince. We solved one bug by posting an "
            "event with the wrong semantic meaning, then had to add another patch to "
            "suppress the wrong-half-of-that-meaning.\n\n"
            "Framing: 'this is what patch-as-last-resort actually costs over time. Each "
            "patch is small. The COMPOSITION of patches gets complicated. Every fake "
            "event needs a guard for the next unintended subscriber.'"
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_mod3_intro(ctx: SlideContext) -> None:
    add_quote_slide(
        ctx.prs,
        quote=JASPER_MIXED_QUOTE,
        attribution=f"{SON_NAME}, feature request #3",
        notes=(
            "The kid's third feature request. After the floods worked, he wanted ANOTHER "
            "mechanic — water that's partly bad, partly clean. Not full badtide, not "
            "drought, something in between.\n\n"
            "This sets up the doing-it-twice section. Walk the audience through what "
            "the second pass at the same patch design revealed."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_contamination_seam(ctx: SlideContext) -> None:
    add_code_slide(
        ctx.prs,
        title="The Mixed Tide seam — contamination IS the mix ratio",
        code=(
            "// Found by decompiling Timberborn.WaterSystem and reading\n"
            "// SimulateContamination:\n"
            "//\n"
            "//   float num3 = (column.Contamination * (num2 + waterDisposed)\n"
            "//                 + contaminationChange) / num;\n"
            "//   _contaminationsBuffer[index3D] =\n"
            "//       (num3 > _maxWaterContamination) ? _maxWaterContamination : num3;\n"
            "//\n"
            "// Translation: WaterSource.Contamination IS the bad-water fraction\n"
            "// emitted at the source. It diffuses into the water-column simulation\n"
            "// as a real mix — not a binary 'is bad / is good' threshold.\n"
            "//\n"
            "// → setting it to 0.3 yields 30% bad / 70% clean. Mechanic is FREE.\n"
            "\n"
            "internal class MixedTideContaminationController : TickableComponent, … {\n"
            "    public override void Tick() {\n"
            "        _waterSourceContamination.SetContamination(_settings.MixedTideContamination);\n"
            "    }\n"
            "}"
        ),
        caption="Read the simulation code, found the variable, set it to the user's chosen value.",
        notes=(
            "The Mixed Tide mechanic is genuinely just 'find the right variable, write "
            "to it on every tick'. The hard part was confirming it was a real ratio and "
            "not a binary threshold downstream.\n\n"
            "Hat-tip to Mechanistry: contamination flows through the water-column "
            "simulation as a real fraction. No hardcoded 'if > 0 emit bad water'. Our "
            "30% setting produces actual 30%/70% mixed water that diffuses correctly."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_demo_mixed_tide(ctx: SlideContext) -> None:
    add_demo_slide(
        ctx.prs,
        title="This is with a mixed tide season",
        steps=[
            "Open Mod Settings — enable Mixed Tide, probability 100%, contamination 30%.",
            "Start a NEW game, fast-forward to the first hazardous cycle.",
            "Point at the date panel — custom 'Mixed Tide' icon + label.",
            "Walk the camera to a water source — partly-bad water emitting + diffusing downstream.",
            "Optional: contrast with a real Badtide cycle (set Mixed Tide off, restart) so the 30%/70% mix vs full bad is visible.",
        ],
        notes=(
            "Pre-prep: a save mid-Mixed-Tide with several water sources visible. The "
            "audience should see the water visibly tinted with contamination, not "
            "fully purple/green like a real badtide. The 30% setting is the headline "
            "mechanic — it's a real ratio that diffuses through the simulation.\n\n"
            "If the contrast-with-real-badtide demo lands, lead with that. Otherwise "
            "just show the mixed tide and move on."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_doing_it_twice(ctx: SlideContext) -> None:
    add_bullets_slide(
        ctx.prs,
        title="Doing it twice — what the second pass revealed",
        bullets=[
            "Every patch from the cascade needed ONE parallel branch for Mixed Tide.",
            ("WeatherUIHelperPatch: drought spec for flood, badtide spec for mixed.", 1),
            ("UIHelperLabelsPatch: 5 × postfix get a MixedTide branch.", 1),
            ("SoundPlayerPatch, NotificationBackgroundPatch, … same shape.", 1),
            "Three new gotchas surfaced along the way:",
            ("`TemplateModule.AddDecorator` needs a paired `Bind<T>().AsTransient()`.", 1),
            ("`IObjectLoader.Get` throws on missing keys — use `loader.Has()` to guard.", 1),
            ("Entity controllers don't auto-wake on save-restore (PostLoad timing).", 1),
        ],
        notes=(
            "The headline insight: the patch DESIGN from Flood Season held up. Mixed "
            "Tide didn't require restructuring anything — every existing patch just "
            "needed one more conditional branch.\n\n"
            "The third gotcha is the architecturally interesting one: when we "
            "force-restore the weather in PostLoad, the entity-side components have "
            "already initialised with the wrong belief about what the current weather "
            "is. They need a wake-up signal — we subscribe to HazardousWeatherSelected"
            "Event guarded by IsHazardousWeather, which fires naturally during our "
            "PostLoad re-emit and stays quiet during vanilla cycle-start."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_gotcha_catalogue(ctx: SlideContext) -> None:
    add_bullets_slide(
        ctx.prs,
        title="Gotchas — nine landmines, documented",
        bullets=[
            "1. bin/obj outside the mod folder (absolute path, not relative)",
            "2. Harmony field-injection needs FOUR underscores for vanilla underscored fields",
            "3. Bindito AsSingleton() is lazy without a tagged interface (ILoadableSingleton)",
            "4. Vanilla type-checks throw on unknown IHazardousWeather (UI helper, sound, fog)",
            "5. Id-keyed lookups are spoofed via FloodWeather.Id = \"DroughtWeather\"",
            "6. Cycle weather is decided at cycle START — settings changes never retro-apply",
            "7. TemplateModule.AddDecorator needs paired Bind<T>().AsTransient()",
            "8. IObjectLoader.Get(PropertyKey<T>) throws on missing keys — use Has()",
            "9. Custom weathers don't auto-fire entity controllers on save-restore",
        ],
        notes=(
            "If anyone takes a photo of one slide during the talk, it'll be this one. "
            "Land it slowly — these are the nine real landmines we hit and wrote down "
            "for next time.\n\n"
            "Each one started as a crash, a wrong-looking screenshot, or a confusing "
            "log line. The discipline that made the list useful: write it down the "
            "moment you've understood the cause, before the context evaporates."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_mod4_rainy_intro(ctx: SlideContext) -> None:
    add_quote_slide(
        ctx.prs,
        quote=JASPER_RAIN_BRIEF,
        attribution=f"{SON_NAME}, feature request #4 (pending)",
        notes=(
            "Closing tease. The mod we haven't built yet.\n\n"
            "Frame it as 'here's the next one — let me walk you through what I'd need "
            "to find to make it work. Note the workflow is the same as the other three: "
            "decompile, find the seam, choose the layer.'"
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_mod4_design(ctx: SlideContext) -> None:
    add_bullets_slide(
        ctx.prs,
        title="Mod #4 design sketch — where would we hook?",
        bullets=[
            "Goal: water visibly falls from the sky and the map's surface gets wet.",
            "Two components: a new IHazardousWeather (familiar) + a way to add water to map tiles.",
            ("Hazard plumbing: copy the FloodWeather pattern. Spoof Id, parallel branches in patches.", 1),
            ("Water-adding mechanic: needs investigation. Candidates worth grepping:", 1),
            ("WaterAdditionService or similar — anything callable per-tile.", 2),
            ("Decomp WaterSystem already for SimulateContamination — grep for AddWater, AddDepth, etc.", 2),
            ("Particle / visual layer: probably orthogonal. Game already has a rain shader for badtide.", 2),
            "Risk: if water-adding isn't exposed, we're back in Harmony-cascade territory.",
        ],
        notes=(
            "This is genuinely the next thing. Don't promise — the talk's already "
            "shipped three mods, the audience doesn't need a fourth.\n\n"
            "Useful talking point: 'this is the workflow you'd use too. Decompile, "
            "grep for keywords matching the behaviour you want, find the seam (or "
            "discover there isn't one), pick the layer, start patching.'\n\n"
            "If anyone asks 'when will this be done?' — Jasper hasn't actually asked "
            "yet, this is me predicting. Don't commit to a timeline."
        ),
        section=ctx.section, slide_num=ctx.slide_num, total=ctx.total,
    )


def _slide_thanks_qa(ctx: SlideContext) -> None:
    add_title_slide(
        ctx.prs,
        title="Thanks",
        subtitle="Questions?",
        sub2=f"github.com/Spycho/timberborn-flood-mod   ·   designed by {SON_NAME}, age {SON_AGE}",
        notes=(
            f"Closing reflection (placeholder, replace with something specific): {CLOSING_REFLECTION}\n\n"
            "Q&A starts here. ~10 minutes.\n\n"
            "Likely questions to prep for:\n"
            "  - 'Why not contribute a PR to Timberborn instead?' — Mechanistry doesn't "
            "take random PRs; the mod ecosystem is the supported extension point.\n"
            "  - 'Is this on the Steam Workshop?' — Not yet. Build is hardened for "
            "publication, just haven't packaged it.\n"
            "  - 'How long did this take?' — Roughly N evenings; specifics in git history.\n"
            "  - 'Could the kid actually read the code?' — He couldn't read most of it. He "
            "understood the SHAPE — 'this part picks the weather', 'this part draws the "
            "picture' — and he made decisions at that level.\n"
            "  - 'How does this relate to our work?' — The decompile / find-the-seam / "
            "patch-when-necessary workflow is what you'd do reverse-engineering any "
            "third-party software in our stack. PlayCanvas asset bundles, vendor SDKs, "
            "anything you don't have source for."
        ),
    )


# --- Driver -----------------------------------------------------------------


def main() -> None:
    out_path = HERE / "timberborn-talk.pptx"
    SCREENSHOTS.mkdir(exist_ok=True)
    prs = build()
    prs.save(out_path)
    print(f"wrote {out_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
